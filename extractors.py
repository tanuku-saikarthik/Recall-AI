import os
import logging
from typing import List, Dict, Any
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from openpyxl import load_workbook

logger = logging.getLogger(__name__)

def extract_pdf(filepath: str) -> List[Dict[str, Any]]:
    """Extracts text from PDF, returning a list of dicts with page numbers."""
    pages = []
    try:
        doc = fitz.open(filepath)
        for i in range(len(doc)):
            page = doc.load_page(i)
            text = page.get_text()
            if text.strip():
                pages.append({"page_or_slide": i + 1, "text": text.strip()})
    except Exception as e:
        logger.error(f"Error extracting PDF {filepath}: {e}")
    return pages

def extract_docx(filepath: str) -> List[Dict[str, Any]]:
    """Extracts text from DOCX, grouping by paragraphs."""
    pages = []
    try:
        doc = docx.Document(filepath)
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(para.text.strip())
        
        if full_text:
            # We don't have pages for docx, so we just return one block or chunk it later
            pages.append({"page_or_slide": 1, "text": "\n".join(full_text)})
    except Exception as e:
        logger.error(f"Error extracting DOCX {filepath}: {e}")
    return pages

def extract_pptx(filepath: str) -> List[Dict[str, Any]]:
    """Extracts text from PPTX, keeping track of slide numbers and notes."""
    slides_data = []
    try:
        prs = Presentation(filepath)
        for i, slide in enumerate(prs.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # Try to get notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_text.append("Notes: " + notes)
                    
            if slide_text:
                slides_data.append({"page_or_slide": i + 1, "text": "\n".join(slide_text)})
    except Exception as e:
        logger.error(f"Error extracting PPTX {filepath}: {e}")
    return slides_data

def extract_xlsx(filepath: str) -> List[Dict[str, Any]]:
    """Extracts text from XLSX, treating each sheet as a 'page'."""
    sheets_data = []
    try:
        wb = load_workbook(filepath, read_only=True, data_only=True)
        for i, sheetname in enumerate(wb.sheetnames):
            sheet = wb[sheetname]
            sheet_text = []
            for row in sheet.iter_rows(values_only=True):
                row_str = " | ".join([str(cell) for cell in row if cell is not None])
                if row_str.strip():
                    sheet_text.append(row_str)
            if sheet_text:
                sheets_data.append({"page_or_slide": i + 1, "text": f"Sheet: {sheetname}\n" + "\n".join(sheet_text)})
    except Exception as e:
        logger.error(f"Error extracting XLSX {filepath}: {e}")
    return sheets_data

def extract_plain_text(filepath: str) -> List[Dict[str, Any]]:
    """Extracts text from plain text and code files."""
    try:
        # Try utf-8 first, fallback to latin-1
        try:
            with open(filepath, 'r', encoding='utf-8') as f:
                content = f.read()
        except UnicodeDecodeError:
            with open(filepath, 'r', encoding='latin-1') as f:
                content = f.read()
                
        if content.strip():
            return [{"page_or_slide": 1, "text": content.strip()}]
    except Exception as e:
        logger.error(f"Error extracting plain text {filepath}: {e}")
    return []

def extract_file_content(filepath: str, extension: str) -> List[Dict[str, Any]]:
    """Routes the file to the appropriate extractor based on extension."""
    ext = extension.lower()
    if ext == '.pdf':
        return extract_pdf(filepath)
    elif ext == '.docx':
        return extract_docx(filepath)
    elif ext == '.pptx':
        return extract_pptx(filepath)
    elif ext == '.xlsx':
        return extract_xlsx(filepath)
    elif ext in ['.txt', '.md', '.py', '.cpp', '.java', '.html', '.json', '.csv', '.js', '.ts', '.css']:
        return extract_plain_text(filepath)
    else:
        return []
